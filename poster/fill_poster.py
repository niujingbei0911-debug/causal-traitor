"""Fill the ICML poster template with the course-project paper content."""

from __future__ import annotations

import os
import subprocess
import sys
from pathlib import Path


def _reexec_with_pptx() -> None:
    try:
        import pptx  # noqa: F401

        return
    except ModuleNotFoundError:
        pass

    if os.environ.get("PPTX_REEXEC") == "1":
        raise

    candidates = [
        os.environ.get("PPTX_PYTHON"),
        "/opt/homebrew/Caskroom/miniconda/base/bin/python",
        "/Users/liyuhang/services/docling/.venv/bin/python",
        "python3",
    ]
    for candidate in [c for c in candidates if c]:
        try:
            probe = subprocess.run(
                [candidate, "-c", "import pptx"],
                stdout=subprocess.DEVNULL,
                stderr=subprocess.DEVNULL,
                check=False,
            )
        except OSError:
            continue
        if probe.returncode == 0:
            env = os.environ.copy()
            env["PPTX_REEXEC"] = "1"
            os.execve(candidate, [candidate, *sys.argv], env)

    raise ModuleNotFoundError("python-pptx is required; install it or set PPTX_PYTHON")


_reexec_with_pptx()

from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Pt


PPTX_PATH = Path("poster/poster_template.pptx")
PIPELINE_PATH = Path("poster/method_pipeline.png")


TITLE = (
    "Selective Adversarial Causal Oversight: A Leakage-Free Benchmark and "
    "Countermodel-Grounded Verifier under Information Asymmetry"
)


TEXT_BLOCKS = {
    "motivation": [
        (
            "LLMs are increasingly used to evaluate causal claims, but real causal "
            "oversight is selective and adversarial under information asymmetry: "
            "claimants hold private information, exploit non-identifiability, and "
            "present persuasive but incomplete reasoning. Existing causal benchmarks "
            "treat unidentifiability as a degraded label and let oracle metadata leak "
            "into the verifier. We formalize the task, release a leakage-free "
            "benchmark, and propose a countermodel-grounded selective verifier."
        )
    ],
    "framing": [
        "Hook: causal claim oversight is selective + adversarial under information asymmetry",
        "Problem: unidentifiable is treated as a degraded label and oracle metadata leaks into verifiers",
        "Method: leakage-free benchmark + four-stage countermodel-grounded selective verifier",
    ],
    "contributions": [
        (
            "C1 (Task) - Selective adversarial causal oversight: three-way verdict "
            "{valid / invalid / unidentifiable} with structured assumption ledger, "
            "refusal reason, and witness slots."
        ),
        (
            "C2 (Benchmark) - Twelve graph families across Pearl L1/L2/L3, "
            "schema-level public/gold partition, five-axis persuasion overlay, "
            "OOD design space."
        ),
        (
            "C3 (Verifier) - Four-stage pipeline (Parser -> Ledger -> Countermodel "
            "Search -> Tool-backed Adjudication) with deterministic decision rule "
            "and abstention gate."
        ),
    ],
    "method_overview": [
        (
            "Four sequential stages turn a natural-language claim into a selective "
            "verdict. Stages 1-2 produce structured slots and an identifying-assumption "
            "ledger; Stage 3 searches for an observationally compatible counter-SCM "
            "that disagrees on the target query; Stage 4 consults Pearl-style "
            "identification tools only when no countermodel survives."
        )
    ],
    "setup_bullets": [
        "Dataset: 12 graph families (4 L1, 5 L2, 3 L3); literature-grounded subset reserved",
        "Splits: train / dev / test_iid / test_ood with five orthogonal holdout axes",
        (
            "Task: three-way verdict + identification status + refusal reason + "
            "assumption ledger + 3 witness slots"
        ),
        "Baselines: Countermodel-Grounded, No-Tools, No-Countermodel, No-Ledger, No-Abstention, Claim-Only",
        "Budget: 3 seeds x 2 samples / family - exploratory",
    ],
    "technical_block": [
        (
            "The third label is not cosmetic. Pearl hierarchy theorem: layer-k data "
            "underdetermines layer-(k+1) truth almost everywhere, so for many claims "
            "the correct answer under public evidence is to abstain with a structured "
            "reason."
        )
    ],
    "proposition": [
        (
            "Proposition 1 (Abstention Necessity). If SCMs M1, M2 both induce public "
            "evidence E but disagree on target query q, then any verifier without an "
            "unidentifiable output is wrong on at least one of M1, M2."
        )
    ],
    "interpretation": [
        "What it does: any yes/no-only verifier is provably wrong on observationally-equivalent samples",
        "Why it is needed: Pearl hierarchy theorem shows layer-k data generically underdetermines layer-(k+1) truth",
        "Design choice: unidentifiable is a first-class label with structured refusal reason",
    ],
    "decision_rule": [
        "Decision rule (fixed five-branch order):",
        "1. strong invalid countermodel -> INVALID",
        "2. compatible models, disagree on q -> UNIDENTIFIABLE (observational_equivalence)",
        "3a. assumption contradicted by tools -> INVALID",
        "3b. ID assumption unsupported -> UNIDENTIFIABLE (missing_identifying_support)",
        "4. tools + ledger jointly support -> VALID",
    ],
    "results_description": [
        (
            "Main benchmark across the full pipeline and four ablations, on the "
            "schema-level public/gold partition. The full system reaches saturated "
            "accuracy on this exploratory budget."
        )
    ],
    "result_caption": [
        "Comparison: pipeline component ablations on same public/gold partition",
        "Best: full system 1.00 / 1.00 (IID / OOD) vs 0.39-0.78 for ablations",
        "Key insight: No-Abstention keeps high raw accuracy but fails on unidentifiable cases",
    ],
    "example_a": [
        (
            "Setup: latent confounding (L1) - observed correlation plus proxy variable; "
            "gold label is unidentifiable"
        ),
        (
            "Output: Stage 3 finds counter-SCM with hidden confounder; verdict "
            "unidentifiable + observational_equivalence refusal reason"
        ),
    ],
    "example_b": [
        (
            "Setup: positive-control verifier with access to one gold-only field on "
            "attack-only benchmark"
        ),
        (
            "Effect: leaking control reaches 1.00 IID accuracy vs 0.33 for clean "
            "public verifier (delta = +0.67)"
        ),
    ],
}


TABLE_ROWS = [
    ["System", "IID", "OOD"],
    ["Countermodel-Grounded", "1.00", "1.00"],
    ["No-Tools", "0.61", "0.61"],
    ["No-Countermodel", "0.61", "0.44"],
    ["No-Abstention", "0.78", "0.94"],
    ["Claim-Only", "0.39", "0.39"],
]


ANCHORS = [
    "Paper Title Goes Here",
    "Author One",
    "Department / Lab Name or Department",
    "Presenter contact",
    "Use this area",
    "Suggested content",
    "Key contributions",
    "Reserve this block",
    "Method figure / pipeline / architecture",
    "Replace with concise setup details or Setup",
    "Setup bullets area",
    "Use this block",
    "Place one key equation here",
    "Interpretation",
    "Additional diagram",
    "Use this section",
    "Primary table or bar chart",
    "Result caption",
    "Example A",
    "Example B",
]


def remove_shape(shape) -> None:
    element = shape._element
    element.getparent().remove(element)


def text_of(shape) -> str:
    if getattr(shape, "has_text_frame", False):
        return shape.text_frame.text or ""
    return ""


def lower_text(shape) -> str:
    return text_of(shape).lower()


def replace_text(
    shape,
    lines,
    *,
    font_size=18,
    bold_first=False,
    clear_if_empty=False,
) -> None:
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    if not lines and clear_if_empty:
        tf.paragraphs[0].text = ""
        return

    for idx, line in enumerate(lines):
        paragraph = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        paragraph.text = line
        paragraph.space_after = Pt(3)
        for run in paragraph.runs:
            run.font.size = Pt(font_size)
            if bold_first and idx == 0:
                run.font.bold = True


def replace_shape_after(shapes, index, lines, *, font_size=18, bold_first=False) -> bool:
    for candidate in shapes[index + 1 :]:
        if getattr(candidate, "has_text_frame", False):
            replace_text(candidate, lines, font_size=font_size, bold_first=bold_first)
            return True
    return False


def set_found(found, anchor, shape_idx, action) -> None:
    found[anchor] = f"shape {shape_idx}: {action}"
    print(f"FOUND: {anchor} -> shape {shape_idx}: {action}")


def insert_results_table(slide, placeholder, container=None) -> None:
    target = container or placeholder
    left, top, width, height = target.left, target.top, target.width, target.height
    remove_shape(placeholder)
    table = slide.shapes.add_table(6, 3, left, top, width, height).table
    table.columns[0].width = int(width * 0.58)
    table.columns[1].width = int(width * 0.21)
    table.columns[2].width = int(width * 0.21)

    for r_idx, row in enumerate(TABLE_ROWS):
        for c_idx, value in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = value
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(15)
                    run.font.bold = r_idx in {0, 1}


def enclosing_container(shapes, target):
    containers = []
    for shape in shapes:
        if shape is target or text_of(shape).strip():
            continue
        if (
            shape.left <= target.left
            and shape.top <= target.top
            and shape.left + shape.width >= target.left + target.width
            and shape.top + shape.height >= target.top + target.height
        ):
            area = shape.width * shape.height
            containers.append((area, shape))
    if not containers:
        return None
    return sorted(containers, key=lambda item: item[0])[0][1]


def fill_slide(slide, found) -> None:
    shapes = list(slide.shapes)
    pending_removals = []

    for index, shape in enumerate(shapes):
        text = lower_text(shape)
        if not text:
            continue

        if "paper title goes here" in text:
            replace_text(shape, [TITLE], font_size=34, bold_first=True)
            set_found(found, "Paper Title Goes Here", index, "title")
        elif "author one" in text:
            replace_text(
                shape,
                ["Authors anonymized for course submission"],
                font_size=23,
            )
            set_found(found, "Author One", index, "authors")
        elif "department / lab name" in text or text.strip().startswith("department"):
            replace_text(
                shape,
                ["Course project - Causal Reasoning - Spring-Summer 2026"],
                font_size=18,
            )
            set_found(found, "Department / Lab Name or Department", index, "affiliation")
        elif "presenter contact" in text:
            replace_text(shape, [], clear_if_empty=True)
            set_found(found, "Presenter contact", index, "cleared")
        elif "use this area" in text:
            replace_text(shape, TEXT_BLOCKS["motivation"], font_size=19)
            set_found(found, "Use this area", index, "motivation")
        elif "suggested content" in text:
            replace_text(shape, ["At a glance"], font_size=16, bold_first=True)
            if replace_shape_after(shapes, index, TEXT_BLOCKS["framing"], font_size=16):
                set_found(found, "Suggested content", index, "framing block")
        elif "key contributions" in text:
            if replace_shape_after(
                shapes, index, TEXT_BLOCKS["contributions"], font_size=15
            ):
                set_found(found, "Key contributions", index, "contribution block")
        elif "reserve this block" in text:
            replace_text(shape, TEXT_BLOCKS["method_overview"], font_size=18)
            set_found(found, "Reserve this block", index, "method overview")
        elif "method figure" in text and ("pipeline" in text or "architecture" in text):
            slide.shapes.add_picture(
                str(PIPELINE_PATH),
                shape.left,
                shape.top,
                width=shape.width,
            )
            pending_removals.append(shape)
            set_found(
                found,
                "Method figure / pipeline / architecture",
                index,
                "embedded method_pipeline.png",
            )
        elif "replace with concise setup details" in text:
            replace_text(shape, ["Setup at a glance"], font_size=16, bold_first=True)
            if replace_shape_after(shapes, index, TEXT_BLOCKS["setup_bullets"], font_size=15):
                set_found(found, "Setup bullets area", index + 1, "setup bullets")
            set_found(
                found,
                "Replace with concise setup details or Setup",
                index,
                "setup header",
            )
        elif "use this block" in text:
            replace_text(shape, TEXT_BLOCKS["technical_block"], font_size=18)
            set_found(found, "Use this block", index, "technical block")
        elif "place one key equation here" in text:
            replace_text(shape, TEXT_BLOCKS["proposition"], font_size=18, bold_first=True)
            set_found(found, "Place one key equation here", index, "proposition")
        elif text.strip() == "interpretation":
            if replace_shape_after(
                shapes, index, TEXT_BLOCKS["interpretation"], font_size=16
            ):
                set_found(found, "Interpretation", index, "interpretation block")
        elif "additional diagram" in text:
            replace_text(
                shape,
                TEXT_BLOCKS["decision_rule"],
                font_size=18,
                bold_first=True,
            )
            set_found(found, "Additional diagram", index, "decision rule")
        elif "use this section" in text:
            replace_text(shape, TEXT_BLOCKS["results_description"], font_size=18)
            set_found(found, "Use this section", index, "results overview")
        elif "primary table or bar chart" in text:
            insert_results_table(slide, shape, enclosing_container(shapes, shape))
            set_found(found, "Primary table or bar chart", index, "results table")
        elif "result caption" in text:
            replace_text(shape, ["Result summary"], font_size=16, bold_first=True)
            if replace_shape_after(
                shapes, index, TEXT_BLOCKS["result_caption"], font_size=16
            ):
                set_found(found, "Result caption", index, "caption block")
        elif text.strip() == "example a":
            if replace_shape_after(shapes, index, TEXT_BLOCKS["example_a"], font_size=14):
                set_found(found, "Example A", index, "example A bullets")
        elif text.strip() == "example b":
            if replace_shape_after(shapes, index, TEXT_BLOCKS["example_b"], font_size=14):
                set_found(found, "Example B", index, "example B bullets")

    for shape in pending_removals:
        remove_shape(shape)


def main() -> None:
    if not PIPELINE_PATH.exists():
        raise FileNotFoundError(f"missing required image: {PIPELINE_PATH}")

    prs = Presentation(PPTX_PATH)
    found = {}
    for slide in prs.slides:
        fill_slide(slide, found)

    prs.save(PPTX_PATH)

    not_found = [anchor for anchor in ANCHORS if anchor not in found]
    total_shapes = sum(len(slide.shapes) for slide in prs.slides)

    print()
    print("SUMMARY")
    print("found:")
    for anchor in ANCHORS:
        if anchor in found:
            print(f"  - {anchor}: {found[anchor]}")
    print("not_found:")
    for anchor in not_found:
        print(f"  - {anchor}")
    print(f"slide_count: {len(prs.slides)}")
    print(f"final_shape_count: {total_shapes}")
    print(f"table_inserted: {'Primary table or bar chart' in found}")

    if not_found:
        print("WARNING: missing anchors:", ", ".join(not_found))


if __name__ == "__main__":
    main()
