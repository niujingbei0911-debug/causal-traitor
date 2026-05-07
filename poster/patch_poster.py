"""Patch the filled poster: fix font sizes and resize oversized text frames."""
from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Pt, Inches

PPTX_PATH = "poster/poster_template.pptx"

# (shape_index, new_font_pt, new_auto_size, new_height_inches_or_None)
# auto_size: 'fit_shape' = TEXT_TO_FIT_SHAPE, 'fit_text' = SHAPE_TO_FIT_TEXT, 'none' = NONE
PATCHES = [
    (11, 20, "fit_text", None),   # motivation: small box → expand
    (14, 18, "fit_text", None),   # at-a-glance framing
    (17, 17, "fit_text", None),   # contributions: 3 long lines
    (21, 20, "fit_text", None),   # method overview desc: 1.14" box → expand
    (28, 18, "fit_text", None),   # setup bullets
    (32, 20, "fit_text", None),   # technical block: 1.1" box → expand
    (34, 20, "fit_text", None),   # proposition: 0.79" box → expand
    (37, 18, "fit_text", None),   # interpretation
    (39, 24, "none",    6.5),     # decision rule: shrink from 15" to 6.5"
    (43, 20, "fit_text", None),   # results description: 1.1" box → expand
    (47, 18, "fit_text", None),   # result caption
    (53, 20, "none",    3.5),     # example A: shrink from 8.18" to 3.5"
    (56, 20, "none",    3.5),     # example B: shrink from 8.18" to 3.5"
]

AUTO_SIZE_MAP = {
    "fit_shape": MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE,
    "fit_text":  MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT,
    "none":      MSO_AUTO_SIZE.NONE,
}


def patch(shape, font_pt, auto_size_key, new_height_in):
    tf = shape.text_frame
    tf.auto_size = AUTO_SIZE_MAP[auto_size_key]
    tf.word_wrap = True
    for para in tf.paragraphs:
        for run in para.runs:
            run.font.size = Pt(font_pt)
    if new_height_in is not None:
        shape.height = Inches(new_height_in)


def main():
    prs = Presentation(PPTX_PATH)
    slide = prs.slides[0]
    shapes = list(slide.shapes)

    for idx, font_pt, auto_size_key, new_h in PATCHES:
        if idx >= len(shapes):
            print(f"WARNING: shape index {idx} out of range")
            continue
        sh = shapes[idx]
        if not sh.has_text_frame:
            print(f"WARNING: shape[{idx}] has no text frame, skipping")
            continue
        old_h = sh.height / 914400
        patch(sh, font_pt, auto_size_key, new_h)
        print(
            f"  [{idx:2d}] font={font_pt}pt auto={auto_size_key} "
            f"height: {old_h:.2f}\" -> {new_h or old_h:.2f}\""
            f" | {sh.text_frame.text[:40].replace(chr(10),' ')!r}"
        )

    prs.save(PPTX_PATH)
    print("\nSaved:", PPTX_PATH)


if __name__ == "__main__":
    main()
