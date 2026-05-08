"""Build the conference-style poster and chart assets from frozen metrics."""

from __future__ import annotations

import json
from pathlib import Path
from typing import Any
from zipfile import ZipFile

import matplotlib.pyplot as plt
import numpy as np
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
POSTER_DIR = ROOT / "poster"
OUTPUTS_DIR = ROOT / "outputs"
PAPER_FACING_DIR = OUTPUTS_DIR / "paper_facing"

FONT = "Aptos"
INK = RGBColor(31, 35, 40)
MUTED = RGBColor(91, 99, 110)
LINE = RGBColor(210, 216, 224)
BG = RGBColor(247, 249, 252)
NAVY = RGBColor(22, 59, 108)
LIGHT_BLUE = RGBColor(234, 242, 255)
LIGHT_GREEN = RGBColor(234, 247, 238)
LIGHT_YELLOW = RGBColor(255, 248, 230)
BLUE = "#2F6FED"
TEAL = "#12805C"
RED = "#C7362F"
GOLD = "#B7791F"
GRAY = "#697386"


def _read_json(path: Path) -> Any:
    return json.loads(path.read_text(encoding="utf-8"))


def _reference_template_path() -> Path | None:
    for path in POSTER_DIR.glob("*.pptx"):
        if path.name != "poster_template.pptx":
            return path
    return None


def _extract_template_logo() -> Path | None:
    template = _reference_template_path()
    if template is None:
        return None
    output = POSTER_DIR / "zju_template_logo.png"
    if output.exists() and output.stat().st_mtime >= template.stat().st_mtime:
        return output
    try:
        with ZipFile(template) as package:
            media = [name for name in package.namelist() if name.startswith("ppt/media/")]
            if not media:
                return None
            image_name = media[0]
            output.write_bytes(package.read(image_name))
    except Exception:
        return None
    return output


def _metric(metrics: dict[str, Any], name: str) -> dict[str, float]:
    item = metrics[name]
    return {
        "mean": float(item["mean"]),
        "ci_lower": float(item.get("ci_lower", item["mean"])),
        "ci_upper": float(item.get("ci_upper", item["mean"])),
        "std": float(item.get("std", 0.0)),
    }


def _setup_text(poster_metrics: dict[str, Any]) -> str:
    setup = poster_metrics["setup"]
    return (
        f"Exploratory snapshot | 3 seeds | "
        f"10/family | diff={setup['difficulty']:.2f} | "
        "95% CI"
    )


def _llm_baseline_text(llm_payload: dict[str, Any]) -> str:
    summary = llm_payload.get("summary", {})
    models = llm_payload.get("models", {})
    best_model = "n/a"
    best_acc = -1.0
    best_unsafe = 0.0
    for model_id, split_payload in models.items():
        ood = split_payload.get("test_ood", {})
        accuracy = float((ood.get("accuracy") or {}).get("mean", -1.0))
        if accuracy > best_acc:
            best_model = "GPT-5.5" if model_id == "gpt55_xhigh" else str(model_id)
            best_acc = accuracy
            best_unsafe = float((ood.get("unsafe_acceptance_rate") or {}).get("mean", 0.0))
    return (
        f"API baselines: 5 models, {int(summary.get('total_predictions', 0))} predictions; "
        f"best OOD {best_model} acc {best_acc:.3f}, unsafe {best_unsafe:.3f}."
    )


def _llm_metric(llm_payload: dict[str, Any], model_id: str, split: str, metric: str) -> float:
    model_payload = llm_payload.get("models", {}).get(model_id, {})
    split_payload = model_payload.get(split, {})
    metric_payload = split_payload.get(metric, {})
    return float(metric_payload.get("mean", 0.0))


def _add_textbox(
    slide: Any,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    *,
    font_size: int = 15,
    bold: bool = False,
    color: RGBColor = INK,
    align: int | None = None,
) -> Any:
    shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.04)
    frame.margin_right = Inches(0.04)
    frame.margin_top = Inches(0.02)
    frame.margin_bottom = Inches(0.02)
    frame.vertical_anchor = MSO_ANCHOR.TOP
    paragraph = frame.paragraphs[0]
    if align is not None:
        paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return shape


def _add_box(
    slide: Any,
    left: float,
    top: float,
    width: float,
    height: float,
    *,
    fill: RGBColor = RGBColor(255, 255, 255),
    line: RGBColor = LINE,
    line_width: float = 0.8,
) -> Any:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    try:
        shape.adjustments[0] = 0.04
    except Exception:
        pass
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(line_width)
    return shape


def _add_template_section(slide: Any, number: str, title: str, left: float, top: float, width: float, height: float) -> None:
    _add_box(slide, left, top, width, height, fill=RGBColor(255, 255, 255), line=LINE, line_width=0.8)
    bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left + 0.19), Inches(top + 0.22), Inches(width - 0.38), Inches(0.61))
    try:
        bar.adjustments[0] = 0.08
    except Exception:
        pass
    bar.fill.solid()
    bar.fill.fore_color.rgb = LIGHT_BLUE
    bar.line.color.rgb = LIGHT_BLUE
    heading = f"{number}. {title}" if number else title
    _add_textbox(slide, heading, left + 0.36, top + 0.27, width - 0.72, 0.44, font_size=25, bold=True, color=NAVY)


def _add_callout(
    slide: Any,
    title: str,
    lines: list[str],
    left: float,
    top: float,
    width: float,
    height: float,
    *,
    fill: RGBColor,
    font_size: int = 16,
) -> None:
    _add_box(slide, left, top, width, height, fill=fill, line=LINE, line_width=0.65)
    _add_textbox(slide, title, left + 0.18, top + 0.14, width - 0.36, 0.35, font_size=18, bold=True, color=MUTED)
    _add_bullets(slide, lines, left + 0.18, top + 0.55, width - 0.36, height - 0.65, font_size=font_size)


def _add_section_title(slide: Any, title: str, left: float, top: float, width: float) -> None:
    _add_textbox(slide, title, left, top, width, 0.25, font_size=16, bold=True, color=RGBColor(20, 83, 136))
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top + 0.28), Inches(width), Inches(0.015))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(20, 83, 136)
    line.line.fill.background()


def _add_panel(slide: Any, left: float, top: float, width: float, height: float) -> None:
    panel = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    panel.fill.solid()
    panel.fill.fore_color.rgb = RGBColor(250, 252, 255)
    panel.line.color.rgb = LINE
    panel.line.width = Pt(0.75)


def _add_chart_frame(slide: Any, image_path: Path, left: float, top: float, width: float, height: float, *, image_width: float | None = None) -> None:
    frame = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    frame.fill.solid()
    frame.fill.fore_color.rgb = RGBColor(255, 255, 255)
    frame.line.color.rgb = RGBColor(91, 141, 239)
    frame.line.width = Pt(0.8)
    pic_width = image_width if image_width is not None else width - 0.7
    slide.shapes.add_picture(str(image_path), Inches(left + (width - pic_width) / 2), Inches(top + 0.2), width=Inches(pic_width))


def _format_bullets(frame: Any, lines: list[str], *, font_size: int = 13, color: RGBColor = INK) -> None:
    frame.clear()
    frame.word_wrap = True
    for idx, line in enumerate(lines):
        paragraph = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        paragraph.text = line
        paragraph.level = 0
        paragraph.font.name = FONT
        paragraph.font.size = Pt(font_size)
        paragraph.font.color.rgb = color
        paragraph.space_after = Pt(4)


def _add_bullets(slide: Any, lines: list[str], left: float, top: float, width: float, height: float, *, font_size: int = 13) -> None:
    shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    shape.text_frame.word_wrap = True
    shape.text_frame.margin_left = Inches(0.05)
    shape.text_frame.margin_right = Inches(0.05)
    shape.text_frame.margin_top = Inches(0.02)
    _format_bullets(shape.text_frame, lines, font_size=font_size)


def _add_flow_step(slide: Any, label: str, left: float, top: float, width: float, height: float, color: RGBColor, *, font_size: int = 10) -> None:
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(250, 252, 255)
    box.line.color.rgb = color
    box.line.width = Pt(1.1)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.06)
    frame.margin_right = Inches(0.06)
    frame.margin_top = Inches(0.03)
    frame.margin_bottom = Inches(0.03)
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = label
    run.font.name = FONT
    run.font.size = Pt(font_size)
    run.font.bold = True
    run.font.color.rgb = INK


def _add_flow_arrow(slide: Any, left: float, top: float) -> None:
    _add_textbox(slide, "->", left, top, 0.28, 0.2, font_size=12, bold=True, color=RGBColor(20, 83, 136), align=PP_ALIGN.CENTER)


def build_charts() -> None:
    poster_metrics = _read_json(PAPER_FACING_DIR / "poster_metrics.json")
    main = poster_metrics["main"]
    ood = poster_metrics["ood"]
    llm_baseline = poster_metrics.get("llm_baseline", {})
    ablation = _read_json(OUTPUTS_DIR / "review_fixed_exp_identifiability_ablation_aggregated_metrics.json")

    plt.rcParams.update({
        "font.family": "DejaVu Sans",
        "axes.spines.top": False,
        "axes.spines.right": False,
        "axes.labelcolor": "#23272f",
        "xtick.color": "#23272f",
        "ytick.color": "#23272f",
    })

    # Main OOD tradeoff chart. Include the strongest API-backed prompt baseline
    # so the poster text and visual comparison refer to the same object.
    systems = [
        ("direct_judge", "Direct"),
        ("tool_baseline", "Tool"),
        ("refusal_aware_baseline", "Refusal-aware"),
        ("gpt55_xhigh", "GPT-5.5"),
        ("countermodel_grounded", "Ours"),
    ]
    acc = []
    acc_low = []
    acc_high = []
    unsafe = []
    labels = []
    for name, label in systems:
        if name == "gpt55_xhigh":
            accuracy = _llm_metric(llm_baseline, name, "test_ood", "accuracy")
            acc.append(accuracy)
            acc_low.append(0.0)
            acc_high.append(0.0)
            unsafe.append(_llm_metric(llm_baseline, name, "test_ood", "unsafe_acceptance_rate"))
        else:
            accuracy = main[name]["test_ood"]["verdict_accuracy"]["mean"]
            acc.append(accuracy)
            acc_low.append(accuracy - main[name]["test_ood"]["verdict_accuracy"]["ci_lower"])
            acc_high.append(main[name]["test_ood"]["verdict_accuracy"]["ci_upper"] - accuracy)
            unsafe.append(main[name]["test_ood"]["unsafe_acceptance_rate"]["mean"])
        labels.append(label)
    fig, ax = plt.subplots(figsize=(7.0, 3.8), dpi=220)
    point_colors = [RED, GOLD, GRAY, "#6B5B95", BLUE]
    sizes = [110, 125, 125, 125, 165]
    ax.errorbar(unsafe, acc, yerr=[acc_low, acc_high], fmt="none", ecolor="#30363D", elinewidth=1.1, capsize=3, zorder=1)
    ax.scatter(unsafe, acc, s=sizes, color=point_colors, alpha=0.95, edgecolor="white", linewidth=1.2, zorder=2)
    offsets = {
        "Direct": (-34, 0),
        "Tool": (8, -12),
        "Refusal-aware": (8, -2),
        "GPT-5.5": (12, 10),
        "Ours": (-40, -2),
    }
    for x_value, y_value, label in zip(unsafe, acc, labels):
        ax.annotate(label, (x_value, y_value), xytext=offsets.get(label, (6, 6)), textcoords="offset points", fontsize=9)
    ax.set_xlim(-0.04, 1.04)
    ax.set_ylim(0, 1.05)
    ax.set_xlabel("Unsafe acceptance rate (lower is better)")
    ax.set_ylabel("OOD accuracy (higher is better)")
    ax.set_title("OOD Accuracy / Safety Tradeoff")
    ax.grid(axis="both", alpha=0.18)
    ax.text(0.04, 0.94, "desired corner", transform=ax.transAxes, ha="left", va="top", fontsize=8.5, color="#145388")
    fig.tight_layout(rect=(0, 0.02, 1, 1))
    fig.savefig(POSTER_DIR / "main_chart.png", transparent=False, facecolor="white")
    plt.close(fig)

    # OOD bucket chart.
    buckets = [
        ("graph_family_ood", "Graph"),
        ("mechanism_ood", "Mechanism"),
        ("attack_family_ood", "Attack"),
        ("context_shift_ood", "Context"),
        ("paired_flip_ood", "Paired flip"),
    ]
    x = np.arange(len(buckets))
    values = [ood[name]["verdict_accuracy"]["mean"] for name, _ in buckets]
    lows = [value - ood[name]["verdict_accuracy"]["ci_lower"] for value, (name, _) in zip(values, buckets)]
    highs = [ood[name]["verdict_accuracy"]["ci_upper"] - value for value, (name, _) in zip(values, buckets)]
    colors = [TEAL, TEAL, GOLD, BLUE, RED]
    fig, ax = plt.subplots(figsize=(7.0, 3.5), dpi=220)
    ax.bar(x, values, color=colors, yerr=[lows, highs], capsize=3)
    ax.set_xticks(x, [label for _, label in buckets], rotation=12, ha="right")
    ax.set_ylim(0, 1.05)
    ax.set_ylabel("Verdict accuracy")
    ax.set_title("OOD Buckets")
    fig.tight_layout(rect=(0, 0.02, 1, 1))
    fig.savefig(POSTER_DIR / "ood_chart.png", transparent=False, facecolor="white")
    plt.close(fig)

    # Ablation tradeoff chart.
    variants = [
        ("countermodel_grounded", "Full"),
        ("no_countermodel", "No CM"),
        ("no_abstention", "No abstain"),
        ("no_tools", "No tools"),
    ]
    recall = [_metric(ablation[name]["test_ood"], "wise_refusal_recall")["mean"] for name, _ in variants]
    over_refusal = [_metric(ablation[name]["test_ood"], "over_refusal_rate")["mean"] for name, _ in variants]
    accuracy = [_metric(ablation[name]["test_ood"], "verdict_accuracy")["mean"] for name, _ in variants]
    fig, ax = plt.subplots(figsize=(7.0, 3.45), dpi=220)
    sizes = [150 + value * 450 for value in accuracy]
    colors = [BLUE, GOLD, RED, GRAY]
    ax.scatter(over_refusal, recall, s=sizes, color=colors, alpha=0.9, edgecolor="white", linewidth=1.2)
    for x_value, y_value, (_, label) in zip(over_refusal, recall, variants):
        ax.annotate(label, (x_value, y_value), xytext=(5, 5), textcoords="offset points", fontsize=9)
    ax.set_xlim(-0.02, max(over_refusal) + 0.10)
    ax.set_ylim(0, 1.05)
    ax.set_xlabel("Over-refusal rate")
    ax.set_ylabel("Wise-refusal recall")
    ax.set_title("Component Ablation (4 Variants)")
    ax.grid(axis="both", alpha=0.18)
    ax.text(0.98, 0.06, "bubble size = OOD accuracy", transform=ax.transAxes, ha="right", va="bottom", fontsize=8.5, color="#5B6370")
    fig.tight_layout(rect=(0, 0.02, 1, 1))
    fig.savefig(POSTER_DIR / "ablation_chart.png", transparent=False, facecolor="white")
    plt.close(fig)


def export_poster_pdf(pptx_path: Path, pdf_path: Path) -> None:
    """Export the poster PPTX to PDF using PowerPoint COM on Windows."""
    try:
        import win32com.client  # type: ignore[import-not-found]
    except ImportError as exc:  # pragma: no cover - platform dependent
        raise RuntimeError("PDF export requires pywin32 and Microsoft PowerPoint on Windows.") from exc

    pptx_path = pptx_path.resolve()
    pdf_path = pdf_path.resolve()
    if pdf_path.exists():
        pdf_path.unlink()

    powerpoint = win32com.client.Dispatch("PowerPoint.Application")
    presentation = None
    try:
        presentation = powerpoint.Presentations.Open(str(pptx_path), True, False, False)
        presentation.SaveAs(str(pdf_path), 32)  # 32 = ppSaveAsPDF
    finally:
        if presentation is not None:
            presentation.Close()
        powerpoint.Quit()


def build_poster() -> None:
    build_charts()
    poster_metrics = _read_json(PAPER_FACING_DIR / "poster_metrics.json")
    manifest = _read_json(PAPER_FACING_DIR / "manifest.json")
    llm_baseline = poster_metrics.get("llm_baseline", {})
    logo_path = _extract_template_logo()

    prs = Presentation()
    prs.slide_width = Inches(37.8)
    prs.slide_height = Inches(28.35)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = BG

    # Template-style frame and title card.
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(37.8), Inches(0.91))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = NAVY
    top_bar.line.fill.background()
    bottom_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(27.83), Inches(37.8), Inches(0.51))
    bottom_bar.fill.solid()
    bottom_bar.fill.fore_color.rgb = NAVY
    bottom_bar.line.fill.background()
    _add_box(slide, 0.87, 0.91, 36.06, 3.07, fill=RGBColor(255, 255, 255), line=LINE)
    _add_textbox(
        slide,
        "Selective Adversarial Causal Oversight",
        1.3,
        1.18,
        28.0,
        0.9,
        font_size=44,
        bold=True,
        color=NAVY,
    )
    _add_textbox(
        slide,
        "Benchmark prototype + countermodel-grounded reference verifier for selective causal refusal",
        1.35,
        2.14,
        27.0,
        0.35,
        font_size=22,
        bold=True,
        color=INK,
    )
    _add_textbox(slide, "Jingbei Niu  •  Zeyuan Tong  •  Yuhang Li", 1.35, 2.6, 19.5, 0.35, font_size=18, bold=True, color=INK)
    _add_textbox(slide, "Zhejiang University, College of Computer Science and Technology", 1.35, 3.0, 20.5, 0.31, font_size=15, color=MUTED)
    _add_textbox(slide, _setup_text(poster_metrics), 1.35, 3.33, 16.0, 0.26, font_size=13, color=RGBColor(32, 96, 180))
    if logo_path is not None:
        slide.shapes.add_picture(str(logo_path), Inches(30.55), Inches(1.34), width=Inches(6.18))

    # 1. Motivation and abstract.
    _add_template_section(slide, "1", "Motivation & Abstract", 0.87, 4.37, 11.65, 8.03)
    _add_textbox(
        slide,
        "Causal oversight systems must decide when a public causal claim is valid and when the available evidence is insufficient. "
        "Current prompt-only and heuristic judges often trade unsafe acceptance against over-refusal. "
        "We study this tension as a benchmark-first problem: the verifier receives only public scenario evidence, while gold SCMs and labels remain evaluator-only. "
        "The current artifact is an exploratory prototype, not a mature benchmark release.",
        1.18,
        5.24,
        11.02,
        1.85,
        font_size=17,
        color=INK,
    )
    _add_callout(
        slide,
        "Key contributions",
        [
            "Selective verdict contract: VALID / INVALID / UNIDENTIFIABLE.",
            "Countermodel-grounded verifier with evaluator-only gold labels.",
            "Frozen evidence package: main tradeoff, API baselines, OOD buckets, and ablations.",
        ],
        1.17,
        7.45,
        11.06,
        1.8,
        fill=LIGHT_GREEN,
        font_size=14,
    )
    _add_callout(
        slide,
        "Task contract",
        [
            "Input: public scenario + causal claim only.",
            "Hidden: SCM, labels, and hidden variables.",
            "Risk metric: unsafe acceptance plus wise refusal.",
        ],
        1.17,
        9.55,
        11.06,
        2.35,
        fill=LIGHT_BLUE,
        font_size=14,
    )

    # 2. Method overview.
    _add_template_section(slide, "2", "Method Overview", 0.87, 12.76, 11.65, 6.93)
    _add_textbox(
        slide,
        "The verifier turns a claim into an assumption ledger, searches for public-compatible counter-SCM witnesses, and uses Pearl-style tools as supporting evidence before issuing a selective verdict.",
        1.18,
        13.62,
        11.02,
        1.0,
        font_size=16,
        color=INK,
    )
    flow_left = 1.2
    step_w = 3.0
    col_gap = 0.55
    row_gap = 0.5
    rows = [
        [
            ("Claim\n+ scenario", RGBColor(183, 121, 31)),
            ("Parser", NAVY),
            ("Assumption\nledger", NAVY),
        ],
        [
            ("Countermodel\nsearch", NAVY),
            ("Tool-backed\nadjudication", NAVY),
            ("Verdict\n+ witness", RGBColor(18, 128, 92)),
        ],
    ]
    for row_idx, row in enumerate(rows):
        top = 15.0 + row_idx * (1.05 + row_gap)
        for col_idx, (label, color) in enumerate(row):
            left = flow_left + col_idx * (step_w + col_gap)
            _add_flow_step(slide, label, left, top, step_w, 1.05, color, font_size=15)
            if col_idx < len(row) - 1:
                _add_flow_arrow(slide, left + step_w + 0.12, top + 0.4)
    _add_textbox(slide, "Counter-SCM witnesses are diagnostic certificates, not leaked gold evidence.", 1.18, 18.35, 11.02, 0.35, font_size=13, color=MUTED)

    # 3. Setup / data / tasks.
    _add_template_section(slide, "3", "Setup / Data / Tasks", 0.87, 20.04, 11.65, 7.24)
    _add_callout(
        slide,
        "Frozen experiment snapshot",
        [
            "Synthetic benchmark families with train/dev/test_iid/test_ood splits.",
            "Three seeds; 10 cases per family; difficulty 0.55.",
            "API baseline matrix: 5 models, 1245 predictions.",
            "All poster numbers are generated from frozen JSON artifacts.",
        ],
        1.17,
        21.04,
        11.06,
        3.05,
        fill=LIGHT_YELLOW,
        font_size=15,
    )
    _add_callout(
        slide,
        "Validity boundary",
        [
            "Prompt-protocol-specific API baselines, not universal LLM claims.",
            "Graph/mechanism OOD buckets are saturated smoke tests.",
            "Release requires larger synthetic runs and real-grounded dual audit.",
        ],
        1.17,
        24.4,
        11.06,
        2.35,
        fill=LIGHT_BLUE,
        font_size=15,
    )

    # 4. Main technical details.
    _add_template_section(slide, "4", "Main Technical Details", 13.07, 4.37, 11.65, 23.41)
    _add_textbox(
        slide,
        "The technical object is selective causal verification: accept identifiable valid/invalid claims, but refuse when public evidence admits multiple compatible causal worlds.",
        13.39,
        5.24,
        11.02,
        1.1,
        font_size=16,
        color=INK,
    )
    _add_box(slide, 13.37, 6.52, 11.06, 1.95, fill=LIGHT_BLUE, line=RGBColor(91, 141, 239), line_width=0.8)
    _add_textbox(
        slide,
        "Selective target: maximize OOD accuracy while minimizing unsafe acceptance and over-refusal.",
        13.62,
        6.83,
        10.55,
        0.55,
        font_size=18,
        bold=True,
        color=RGBColor(32, 96, 180),
        align=PP_ALIGN.CENTER,
    )
    _add_textbox(slide, "Key mechanism: public-compatible countermodels witness missing identification assumptions.", 13.65, 7.45, 10.5, 0.35, font_size=13, color=MUTED, align=PP_ALIGN.CENTER)
    _add_callout(
        slide,
        "Interpretation",
        [
            "A countermodel witness supports UNIDENTIFIABLE rather than forcing binary validity.",
            "Pearl-style tools constrain the verifier but do not expose evaluator-only labels.",
            "The design targets unsafe acceptance without collapsing into blanket refusal.",
        ],
        13.37,
        8.75,
        11.06,
        3.0,
        fill=LIGHT_GREEN,
        font_size=14,
    )
    _add_chart_frame(slide, POSTER_DIR / "ood_chart.png", 13.37, 12.05, 11.06, 5.15, image_width=9.95)
    _add_callout(
        slide,
        "OOD stress reading",
        [
            "Paired-flip is the clearest current failure slice: acc 0.439.",
            "Graph/mechanism buckets are saturated smoke tests.",
            "Use these buckets as diagnostics, not final robustness proof.",
        ],
        13.37,
        17.65,
        11.06,
        2.4,
        fill=LIGHT_YELLOW,
        font_size=15,
    )
    _add_chart_frame(slide, POSTER_DIR / "ablation_chart.png", 13.37, 20.55, 11.06, 5.65, image_width=10.05)
    _add_textbox(slide, "Ablation uses four frozen component variants; bubble size encodes OOD accuracy.", 13.7, 26.35, 10.4, 0.35, font_size=12, color=MUTED, align=PP_ALIGN.CENTER)

    # 5. Main results and examples.
    _add_template_section(slide, "5", "Main Results / Examples", 25.28, 4.37, 11.65, 9.72)
    _add_textbox(
        slide,
        "The strongest quantitative evidence is the selective OOD tradeoff, followed by bucket-level stress tests and component ablations.",
        25.59,
        5.24,
        11.02,
        0.95,
        font_size=16,
        color=INK,
    )
    _add_chart_frame(slide, POSTER_DIR / "main_chart.png", 25.57, 6.15, 11.06, 5.1, image_width=9.75)
    _add_callout(
        slide,
        "Result caption",
        [
            "Ours: OOD acc 0.891, unsafe accept 0.000.",
            "GPT-5.5 prompt baseline: OOD acc 0.545, unsafe accept 0.071.",
            "The evidence supports a prototype diagnostic claim, not a final benchmark release.",
        ],
        25.57,
        11.55,
        11.06,
        2.15,
        fill=LIGHT_YELLOW,
        font_size=14,
    )

    # Examples and gates.
    _add_template_section(slide, "", "Examples", 25.28, 14.57, 11.65, 13.13)
    _add_box(slide, 25.59, 15.72, 5.43, 6.35, fill=LIGHT_BLUE, line=LINE)
    _add_textbox(slide, "Example A: unidentifiable claim", 25.75, 15.88, 5.09, 0.35, font_size=15, bold=True, color=MUTED)
    _add_textbox(
        slide,
        "Claim: association between therapy_flag and yield_score is treated as causal.\n\nGold label: UNIDENTIFIABLE.\n\nReason: public evidence leaves confounding or reverse causality open.",
        25.75,
        16.45,
        5.05,
        4.9,
        font_size=13,
        color=INK,
    )
    _add_box(slide, 31.2, 15.72, 5.43, 6.35, fill=LIGHT_GREEN, line=LINE)
    _add_textbox(slide, "Example B: verifier response", 31.37, 15.88, 5.09, 0.35, font_size=15, bold=True, color=MUTED)
    _add_textbox(
        slide,
        "Output: refuse the causal conclusion.\n\nWitness: construct a public-compatible counter-SCM.\n\nTakeaway: refusal is evidence-sensitive, not blanket abstention.",
        31.37,
        16.45,
        5.05,
        4.9,
        font_size=13,
        color=INK,
    )
    _add_callout(
        slide,
        "Validity gates before release",
        [
            "Leakage check: oracle-positive 1.000/1.000; gain is contamination.",
            _llm_baseline_text(llm_baseline),
            "Next gates: larger synthetic runs, real-grounded dual audit, frozen model versions.",
        ],
        25.57,
        22.55,
        11.06,
        3.15,
        fill=LIGHT_YELLOW,
        font_size=14,
    )
    _add_textbox(
        slide,
        f"Frozen artifact manifest: outputs/paper_facing/manifest.json | {manifest['api_model_baseline_status']}",
        25.6,
        26.45,
        10.9,
        0.32,
        font_size=11,
        color=MUTED,
        align=PP_ALIGN.CENTER,
    )

    pptx_path = POSTER_DIR / "poster_template.pptx"
    pdf_path = POSTER_DIR / "poster_template.pdf"
    prs.save(pptx_path)
    export_poster_pdf(pptx_path, pdf_path)


if __name__ == "__main__":
    build_poster()
